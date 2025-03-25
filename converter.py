import os
import logging
from pathlib import Path
from typing import Callable, Dict
import pandas as pd
import numpy as np
from pdf2docx import Converter
from docx import Document
from openpyxl import load_workbook, Workbook
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import markdown
import csv
import tempfile
import fitz
from PyPDF2 import PdfReader
from pptx import Presentation
from docx2pdf import convert
import mammoth
import html2text
import pdfplumber

# ver 13.01 -> nâng cấp chuyển đổi từ pdf->docx
import io
import re
import pytesseract
import shutil
import subprocess
from PIL import Image
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Thiết lập logging
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)


def chuyen_doi_an_toan(func: Callable) -> Callable:
    """Decorator để xử lý ngoại lệ trong các hàm chuyển đổi."""

    def wrapper(*args, **kwargs):
        try:
            return func(*args, **kwargs)
        except Exception as e:
            logging.error(f"Lỗi trong {func.__name__}: {str(e)}")
            return None

    return wrapper


@chuyen_doi_an_toan
def chuyen_doi_pdf_sang_docx(duong_dan_pdf: str) -> str:  # đã được nâng cấp 25/03
    """
    Chuyển đổi PDF sang DOCX với khả năng bảo toàn công thức toán học và hình ảnh
    sử dụng phương pháp hoàn toàn miễn phí.

    Args:
        duong_dan_pdf: Đường dẫn đến tệp PDF

    Returns:
        Đường dẫn đến tệp DOCX đã tạo
    """
    duong_dan_docx = Path(duong_dan_pdf).with_suffix(".docx")

    # Kiểm tra xem LibreOffice có sẵn không (phương pháp tốt nhất)
    if has_libreoffice():
        try:
            logging.info("Đang thử chuyển đổi bằng LibreOffice...")
            result = convert_with_libreoffice(duong_dan_pdf, duong_dan_docx)
            if result:
                logging.info(
                    f"Đã chuyển đổi thành công bằng LibreOffice: {duong_dan_docx}"
                )
                return str(duong_dan_docx)
        except Exception as e:
            logging.error(f"Lỗi khi sử dụng LibreOffice: {str(e)}")

    # Phương pháp 2: Sử dụng phương pháp hình ảnh + văn bản
    try:
        # Tạo document mới
        doc = Document()

        # Thiết lập font và kích thước mặc định
        style = doc.styles["Normal"]
        style.font.name = "Times New Roman"
        style.font.size = Pt(12)

        # Mở PDF bằng PyMuPDF
        pdf_doc = fitz.open(duong_dan_pdf)
        total_pages = len(pdf_doc)

        # Tạo thư mục tạm
        with tempfile.TemporaryDirectory() as temp_dir:
            logging.info(f"Đang xử lý PDF có {total_pages} trang...")

            for page_num, page in enumerate(pdf_doc):
                logging.info(f"Đang xử lý trang {page_num + 1}/{total_pages}")

                # Thêm tiêu đề trang
                heading = doc.add_heading(f"Trang {page_num + 1}", level=1)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

                # Phân tích trang để tìm các vùng văn bản và công thức
                blocks = page.get_text("dict")["blocks"]

                # Trích xuất toàn bộ trang dưới dạng hình ảnh với độ phân giải cao
                page_img = page.get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72))
                page_img_path = os.path.join(temp_dir, f"page_{page_num + 1}.png")
                page_img.save(page_img_path)

                # Phân tích các khối để xác định văn bản thông thường và công thức
                for block_idx, block in enumerate(blocks):
                    if block["type"] == 0:  # Khối văn bản
                        block_text = ""
                        for line in block["lines"]:
                            line_text = ""
                            for span in line["spans"]:
                                line_text += span["text"] + " "
                            block_text += line_text.strip() + "\n"

                        # Kiểm tra xem có phải công thức không
                        if is_likely_formula(block_text):
                            # Trích xuất công thức dưới dạng hình ảnh
                            formula_rect = fitz.Rect(block["bbox"])
                            formula_img = page.get_pixmap(
                                clip=formula_rect, matrix=fitz.Matrix(3, 3)
                            )

                            formula_img_path = os.path.join(
                                temp_dir, f"page_{page_num + 1}_formula_{block_idx}.png"
                            )
                            formula_img.save(formula_img_path)

                            # Thêm hình ảnh công thức vào tài liệu
                            para = doc.add_paragraph()
                            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            run = para.add_run()
                            run.add_picture(formula_img_path, width=None, height=None)
                        else:
                            # Thêm văn bản thông thường
                            if block_text.strip():
                                para = doc.add_paragraph(block_text.strip())

                    elif block["type"] == 1:  # Khối hình ảnh
                        # Trích xuất hình ảnh từ trang
                        img_rect = fitz.Rect(block["bbox"])
                        img = page.get_pixmap(clip=img_rect, matrix=fitz.Matrix(2, 2))

                        img_path = os.path.join(
                            temp_dir, f"page_{page_num + 1}_block_img_{block_idx}.png"
                        )
                        img.save(img_path)

                        # Thêm hình ảnh vào tài liệu
                        para = doc.add_paragraph()
                        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run = para.add_run()
                        run.add_picture(img_path, width=None, height=None)

                # Trích xuất các hình ảnh nhúng
                image_list = page.get_images(full=True)
                for img_idx, img_info in enumerate(image_list):
                    xref = img_info[0]
                    base_image = pdf_doc.extract_image(xref)
                    image_bytes = base_image["image"]

                    # Lưu hình ảnh
                    img_filename = os.path.join(
                        temp_dir, f"page_{page_num + 1}_img_{img_idx}.png"
                    )
                    with open(img_filename, "wb") as img_file:
                        img_file.write(image_bytes)

                    # Kiểm tra kích thước hình ảnh
                    try:
                        pil_img = Image.open(img_filename)
                        if (
                            pil_img.width > 50 and pil_img.height > 50
                        ):  # Bỏ qua hình ảnh quá nhỏ
                            para = doc.add_paragraph()
                            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            run = para.add_run()
                            run.add_picture(img_filename, width=None, height=None)
                    except Exception as img_err:
                        logging.warning(f"Không thể xử lý hình ảnh: {str(img_err)}")

                # Thêm ngắt trang sau mỗi trang PDF
                if page_num < total_pages - 1:
                    doc.add_page_break()

            # Lưu tài liệu
            doc.save(duong_dan_docx)
            logging.info(f"Đã lưu tài liệu DOCX: {duong_dan_docx}")

            return str(duong_dan_docx)

    except Exception as e:
        logging.error(f"Lỗi khi chuyển đổi PDF sang DOCX: {str(e)}")

        # Phương pháp dự phòng: Chuyển toàn bộ trang thành hình ảnh
        try:
            logging.info("Đang thử phương pháp dự phòng...")

            # Tạo document mới
            doc = Document()

            # Mở PDF bằng PyMuPDF
            pdf_doc = fitz.open(duong_dan_pdf)

            for page_num, page in enumerate(pdf_doc):
                # Thêm tiêu đề trang
                heading = doc.add_heading(f"Trang {page_num + 1}", level=1)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

                # Chuyển trang thành hình ảnh với độ phân giải cao
                pix = page.get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72))
                img_bytes = pix.tobytes("png")

                # Thêm hình ảnh trang vào tài liệu
                para = doc.add_paragraph()
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = para.add_run()
                run.add_picture(io.BytesIO(img_bytes), width=Inches(6))

                # Thêm ngắt trang
                if page_num < len(pdf_doc) - 1:
                    doc.add_page_break()

            # Lưu tài liệu
            doc.save(duong_dan_docx)
            logging.info(
                f"Đã lưu tài liệu DOCX (phương pháp dự phòng): {duong_dan_docx}"
            )

            return str(duong_dan_docx)

        except Exception as e2:
            logging.error(f"Phương pháp dự phòng cũng thất bại: {str(e2)}")
            return None


def is_likely_formula(text):  # đã được nâng cấp 25/03
    """Kiểm tra xem một đoạn văn bản có khả năng là công thức toán học không."""
    # Các mẫu regex để nhận dạng công thức toán học
    math_patterns = [
        r"[=+\-*/^]",  # Các toán tử cơ bản
        r"\\[a-zA-Z]+",  # Lệnh LaTeX
        r"[α-ωΑ-Ω]",  # Ký tự Hy Lạp
        r"[∫∑∏√∂∇∆]",  # Ký hiệu toán học
        r"\$.*?\$",  # Công thức LaTeX
        r"$$\d+$$",  # Tham chiếu phương trình
        r"_{.*?}",  # Chỉ số dưới
        r"\^{.*?}",  # Chỉ số trên
    ]

    # Kiểm tra các mẫu công thức
    for pattern in math_patterns:
        if re.search(pattern, text) and not pattern == r"[=+\-*/^]":
            return True
        elif pattern == r"[=+\-*/^]" and len(re.findall(pattern, text)) > 3:
            return True

    # Kiểm tra tỷ lệ ký tự đặc biệt
    special_chars = sum(1 for c in text if c in "=+-*/^()[]{}\\<>~_|")
    if len(text) > 0 and special_chars / len(text) > 0.15:
        return True

    return False


def has_libreoffice():  # đã được nâng cấp 25/03 - cần cài sẵn libreoffice
    """Kiểm tra xem LibreOffice có được cài đặt không."""
    try:
        # Kiểm tra lệnh soffice (LibreOffice)
        subprocess.run(["soffice", "--version"], capture_output=True, check=False)
        return True
    except (FileNotFoundError, subprocess.SubprocessError):
        # Kiểm tra lệnh libreoffice
        try:
            subprocess.run(
                ["libreoffice", "--version"], capture_output=True, check=False
            )
            return True
        except (FileNotFoundError, subprocess.SubprocessError):
            return False


def convert_with_libreoffice(input_path, output_path):  # đã được nâng cấp 25/03
    """Sử dụng LibreOffice để chuyển đổi PDF sang DOCX."""
    try:
        # Tạo thư mục tạm
        with tempfile.TemporaryDirectory() as temp_dir:
            # Xác định lệnh LibreOffice
            libreoffice_cmd = "soffice"
            try:
                subprocess.run(
                    [libreoffice_cmd, "--version"], capture_output=True, check=False
                )
            except (FileNotFoundError, subprocess.SubprocessError):
                libreoffice_cmd = "libreoffice"

            # Lệnh chuyển đổi
            cmd = [
                libreoffice_cmd,
                "--headless",
                "--convert-to",
                "docx",
                "--outdir",
                temp_dir,
                str(input_path),
            ]

            # Thực thi lệnh
            process = subprocess.run(cmd, capture_output=True, text=True, check=False)

            if process.returncode == 0:
                # Tìm tệp đầu ra
                output_filename = Path(input_path).stem + ".docx"
                temp_output_path = os.path.join(temp_dir, output_filename)

                if os.path.exists(temp_output_path):
                    # Sao chép tệp kết quả
                    shutil.copy2(temp_output_path, output_path)
                    return True

            logging.error(f"LibreOffice không thành công: {process.stderr}")
            return False

    except Exception as e:
        logging.error(f"Lỗi khi sử dụng LibreOffice: {str(e)}")
        return False


def check_for_math_formulas(pdf_path):  # đã được nâng cấp 25/03
    """Kiểm tra xem PDF có chứa công thức toán học không."""
    try:
        # Mở PDF bằng PyMuPDF
        doc = fitz.open(pdf_path)

        # Các mẫu regex để nhận dạng công thức toán học
        math_patterns = [
            r"[=+\-*/^]",  # Các toán tử cơ bản
            r"\\[a-zA-Z]+",  # Lệnh LaTeX
            r"[α-ωΑ-Ω]",  # Ký tự Hy Lạp
            r"[∫∑∏√∂∇∆]",  # Ký hiệu toán học
            r"\$.*?\$",  # Công thức LaTeX
            r"$$\d+$$",  # Tham chiếu phương trình
            r"_{.*?}",  # Chỉ số dưới
            r"\^{.*?}",  # Chỉ số trên
        ]

        # Kiểm tra một số trang đầu tiên
        pages_to_check = min(5, len(doc))
        for page_num in range(pages_to_check):
            page = doc[page_num]
            text = page.get_text()

            # Kiểm tra các mẫu công thức
            for pattern in math_patterns:
                if re.search(pattern, text) and not pattern == r"[=+\-*/^]":
                    return True
                elif pattern == r"[=+\-*/^]" and len(re.findall(pattern, text)) > 10:
                    return True

        return False
    except Exception as e:
        logging.warning(f"Không thể kiểm tra công thức toán học: {str(e)}")
        return False


def is_likely_formula(text):  # đã được nâng cấp 25/03
    """Kiểm tra xem một đoạn văn bản có khả năng là công thức toán học không."""
    # Các mẫu regex để nhận dạng công thức toán học
    math_patterns = [
        r"[=+\-*/^]",  # Các toán tử cơ bản
        r"\\[a-zA-Z]+",  # Lệnh LaTeX
        r"[α-ωΑ-Ω]",  # Ký tự Hy Lạp
        r"[∫∑∏√∂∇∆]",  # Ký hiệu toán học
        r"\$.*?\$",  # Công thức LaTeX
        r"$$\d+$$",  # Tham chiếu phương trình
        r"_{.*?}",  # Chỉ số dưới
        r"\^{.*?}",  # Chỉ số trên
    ]

    # Kiểm tra các mẫu công thức
    for pattern in math_patterns:
        if re.search(pattern, text) and not pattern == r"[=+\-*/^]":
            return True
        elif pattern == r"[=+\-*/^]" and len(re.findall(pattern, text)) > 3:
            return True

    # Kiểm tra tỷ lệ ký tự đặc biệt
    special_chars = sum(1 for c in text if c in "=+-*/^()[]{}\\<>~_|")
    if len(text) > 0 and special_chars / len(text) > 0.2:
        return True

    return False


def verify_conversion_quality(pdf_path, docx_path):  # đã được nâng cấp 25/03
    """Kiểm tra chất lượng chuyển đổi bằng cách so sánh số lượng văn bản."""
    try:
        # Đọc văn bản từ PDF
        pdf_doc = fitz.open(pdf_path)
        pdf_text = ""
        for page in pdf_doc:
            pdf_text += page.get_text()

        # Đọc văn bản từ DOCX
        doc = Document(docx_path)
        docx_text = ""
        for para in doc.paragraphs:
            docx_text += para.text + "\n"

        # So sánh độ dài văn bản (loại bỏ khoảng trắng)
        pdf_text_len = len(re.sub(r"\s+", "", pdf_text))
        docx_text_len = len(re.sub(r"\s+", "", docx_text))

        # Nếu DOCX chứa ít nhất 70% văn bản từ PDF, coi là chất lượng tốt
        if pdf_text_len > 0 and docx_text_len / pdf_text_len >= 0.7:
            return True
        return False
    except Exception as e:
        logging.warning(f"Không thể kiểm tra chất lượng chuyển đổi: {str(e)}")
        return True  # Mặc định là chấp nhận


def has_libreoffice():  # đã được nâng cấp 25/03
    """Kiểm tra xem LibreOffice có được cài đặt không."""
    try:
        subprocess.run(["soffice", "--version"], capture_output=True)
        return True
    except (FileNotFoundError, subprocess.SubprocessError):
        return False


# chuyển đổi từ pdf-docx vừa được cải tiến, ver13.01 - 25/03
@chuyen_doi_an_toan
def chuyen_doi_pdf_sang_xlsx(duong_dan_pdf: str) -> str:
    """
    Chuyển đổi PDF sang XLSX với khả năng phát hiện và trích xuất bảng.

    Args:
        duong_dan_pdf: Đường dẫn đến tệp PDF

    Returns:
        Đường dẫn đến tệp XLSX đã tạo
    """
    duong_dan_xlsx = Path(duong_dan_pdf).with_suffix(".xlsx")

    try:
        with pdfplumber.open(duong_dan_pdf) as pdf:

            with pd.ExcelWriter(duong_dan_xlsx, engine="openpyxl") as writer:
                for i, page in enumerate(pdf.pages):

                    tables = page.extract_tables()

                    if tables:

                        for j, table in enumerate(tables):
                            df = pd.DataFrame(
                                table[1:], columns=table[0] if table else None
                            )
                            sheet_name = f"Trang_{i+1}_Bảng_{j+1}"
                            df.to_excel(writer, sheet_name=sheet_name, index=False)
                            logging.info(f"Đã trích xuất bảng {j+1} từ trang {i+1}")
                    else:

                        text = page.extract_text()
                        if text:
                            lines = [
                                line.split()
                                for line in text.split("\n")
                                if line.strip()
                            ]
                            df = pd.DataFrame(lines)
                            sheet_name = f"Trang_{i+1}_Văn_bản"
                            df.to_excel(
                                writer, sheet_name=sheet_name, index=False, header=False
                            )

        logging.info(
            f"Đã chuyển đổi thành công '{duong_dan_pdf}' sang '{duong_dan_xlsx}'"
        )
        return str(duong_dan_xlsx)

    except Exception as e:
        logging.error(f"Lỗi khi chuyển đổi PDF sang XLSX: {str(e)}")
        return None


@chuyen_doi_an_toan
def chuyen_doi_docx_sang_pdf(duong_dan_docx: str) -> str:
    duong_dan_pdf = Path(duong_dan_docx).with_suffix(".pdf")
    convert(duong_dan_docx, str(duong_dan_pdf))
    return str(duong_dan_pdf)


@chuyen_doi_an_toan
def chuyen_doi_xlsx_sang_docx(duong_dan_xlsx: str) -> str:
    duong_dan_docx = Path(duong_dan_xlsx).with_suffix(".docx")
    df = pd.read_excel(duong_dan_xlsx)

    doc = Document()
    for column in df.columns:
        doc.add_heading(column, level=1)
        for value in df[column]:
            doc.add_paragraph(str(value))
        doc.add_paragraph()  # Thêm một dòng trống giữa các cột

    doc.save(duong_dan_docx)
    return str(duong_dan_docx)


@chuyen_doi_an_toan
def chuyen_doi_docx_sang_xlsx(duong_dan_docx: str) -> str:
    duong_dan_xlsx = Path(duong_dan_docx).with_suffix(".xlsx")
    doc = Document(duong_dan_docx)

    data = []
    for paragraph in doc.paragraphs:
        data.append([paragraph.text])

    df = pd.DataFrame(data)
    df.to_excel(duong_dan_xlsx, index=False, header=False)
    return str(duong_dan_xlsx)


@chuyen_doi_an_toan
def chuyen_doi_xlsx_sang_pdf(duong_dan_xlsx: str) -> str:
    duong_dan_pdf = Path(duong_dan_xlsx).with_suffix(".pdf")
    df = pd.read_excel(duong_dan_xlsx)

    pdf = canvas.Canvas(str(duong_dan_pdf), pagesize=letter)
    y = 750  # Tọa độ y bắt đầu
    for column in df.columns:
        pdf.drawString(100, y, column)
        y -= 20
        for value in df[column]:
            pdf.drawString(120, y, str(value))
            y -= 15
            if y < 50:  # Bắt đầu trang mới nếu gần đến cuối trang
                pdf.showPage()
                y = 750
    pdf.save()

    return str(duong_dan_pdf)


@chuyen_doi_an_toan
def chuyen_doi_xlsx_sang_csv(duong_dan_xlsx: str) -> str:
    duong_dan_csv = Path(duong_dan_xlsx).with_suffix(".csv")
    df = pd.read_excel(duong_dan_xlsx)
    df.to_csv(duong_dan_csv, index=False)
    return str(duong_dan_csv)


@chuyen_doi_an_toan
def chuyen_doi_txt_sang_docx(duong_dan_txt: str) -> str:
    duong_dan_docx = Path(duong_dan_txt).with_suffix(".docx")
    with open(duong_dan_txt, "r", encoding="utf-8") as file:
        text = file.read()

    doc = Document()
    for paragraph in text.split("\n"):
        doc.add_paragraph(paragraph)
    doc.save(duong_dan_docx)

    return str(duong_dan_docx)


@chuyen_doi_an_toan
def chuyen_doi_txt_sang_pdf(duong_dan_txt: str) -> str:
    duong_dan_pdf = Path(duong_dan_txt).with_suffix(".pdf")
    with open(duong_dan_txt, "r", encoding="utf-8") as file:
        text = file.read()

    pdf = canvas.Canvas(str(duong_dan_pdf), pagesize=letter)
    y = 750  # Tọa độ y bắt đầu
    for line in text.split("\n"):
        pdf.drawString(100, y, line)
        y -= 15
        if y < 50:  # Bắt đầu trang mới nếu gần đến cuối trang
            pdf.showPage()
            y = 750
    pdf.save()

    return str(duong_dan_pdf)


@chuyen_doi_an_toan
def chuyen_doi_txt_sang_md(duong_dan_txt: str) -> str:
    duong_dan_md = Path(duong_dan_txt).with_suffix(".md")

    with open(duong_dan_txt, "r", encoding="utf-8") as file:
        noi_dung_txt = file.read()

    # Chuyển đổi đơn giản: giả định mỗi dòng là một đoạn văn
    noi_dung_md = "\n\n".join(noi_dung_txt.split("\n"))

    with open(duong_dan_md, "w", encoding="utf-8") as file:
        file.write(noi_dung_md)

    return str(duong_dan_md)


@chuyen_doi_an_toan
def chuyen_doi_csv_sang_xlsx(duong_dan_csv: str) -> str:
    duong_dan_xlsx = Path(duong_dan_csv).with_suffix(".xlsx")
    df = pd.read_csv(duong_dan_csv)
    df.to_excel(duong_dan_xlsx, index=False)
    return str(duong_dan_xlsx)


@chuyen_doi_an_toan
def chuyen_doi_pptx_sang_pdf(duong_dan_pptx: str) -> str:
    from win32com import client

    duong_dan_pdf = Path(duong_dan_pptx).with_suffix(".pdf")

    powerpoint = client.Dispatch("Powerpoint.Application")
    deck = powerpoint.Presentations.Open(duong_dan_pptx)
    deck.SaveAs(duong_dan_pdf, 32)  # 32 là mã định dạng PDF
    deck.Close()
    powerpoint.Quit()

    return str(duong_dan_pdf)


@chuyen_doi_an_toan
def chuyen_doi_pptx_sang_docx(duong_dan_pptx: str) -> str:
    duong_dan_docx = Path(duong_dan_pptx).with_suffix(".docx")

    presentation = Presentation(duong_dan_pptx)
    doc = Document()

    for slide in presentation.slides:
        if slide.shapes.title:
            doc.add_heading(slide.shapes.title.text, level=1)
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                doc.add_paragraph(shape.text)
        doc.add_page_break()

    doc.save(duong_dan_docx)
    return str(duong_dan_docx)


@chuyen_doi_an_toan
def chuyen_doi_md_sang_txt(duong_dan_md: str) -> str:
    duong_dan_txt = Path(duong_dan_md).with_suffix(".txt")

    with open(duong_dan_md, "r", encoding="utf-8") as file:
        noi_dung_md = file.read()

    h = html2text.HTML2Text()
    h.ignore_links = True
    noi_dung_txt = h.handle(markdown.markdown(noi_dung_md))

    with open(duong_dan_txt, "w", encoding="utf-8") as file:
        file.write(noi_dung_txt)

    return str(duong_dan_txt)


@chuyen_doi_an_toan
def chuyen_doi_md_sang_html(duong_dan_md: str) -> str:
    duong_dan_html = Path(duong_dan_md).with_suffix(".html")

    with open(duong_dan_md, "r", encoding="utf-8") as file:
        noi_dung_md = file.read()

    noi_dung_html = markdown.markdown(noi_dung_md, extensions=["extra"])

    with open(duong_dan_html, "w", encoding="utf-8") as file:
        file.write(f"<html><body>{noi_dung_html}</body></html>")

    return str(duong_dan_html)


@chuyen_doi_an_toan
def chuyen_doi_pdf_sang_txt(duong_dan_pdf: str) -> str:
    """
    Chuyển đổi PDF sang TXT với hỗ trợ OCR cho PDF quét và cải thiện định dạng.

    Args:
        duong_dan_pdf: Đường dẫn đến tệp PDF

    Returns:
        Đường dẫn đến tệp TXT đã tạo
    """
    duong_dan_txt = Path(duong_dan_pdf).with_suffix(".txt")

    # Phương pháp 1: Sử dụng pdfplumber (tốt cho PDF có văn bản)
    try:
        logging.info(f"Đang chuyển đổi '{duong_dan_pdf}' sang TXT bằng pdfplumber...")

        with pdfplumber.open(duong_dan_pdf) as pdf:
            with open(duong_dan_txt, "w", encoding="utf-8") as txt_file:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text:
                        txt_file.write(f"--- Trang {i+1} ---\n\n")
                        txt_file.write(text)
                        txt_file.write("\n\n")

        # Kiểm tra kết quả
        with open(duong_dan_txt, "r", encoding="utf-8") as f:
            content = f.read().strip()

        if len(content) > 100:  # Nếu có đủ văn bản
            logging.info(f"Đã chuyển đổi thành công bằng pdfplumber: '{duong_dan_txt}'")
            return str(duong_dan_txt)
        else:
            logging.warning(
                "pdfplumber trích xuất ít văn bản. Có thể là PDF quét. Thử OCR..."
            )
    except Exception as e:
        logging.warning(f"pdfplumber không thành công: {str(e)}")

    # Phương pháp 2: Sử dụng PyMuPDF + OCR cho PDF quét
    try:
        logging.info("Đang thử chuyển đổi bằng PyMuPDF + OCR...")

        # Tạo thư mục tạm cho hình ảnh
        with tempfile.TemporaryDirectory() as temp_dir:
            # Mở PDF
            pdf_doc = fitz.open(duong_dan_pdf)

            with open(duong_dan_txt, "w", encoding="utf-8") as txt_file:
                for page_num, page in enumerate(pdf_doc):
                    # Thử trích xuất văn bản trực tiếp
                    text = page.get_text("text")

                    # Nếu ít văn bản, thử OCR
                    if len(text.strip()) < 100:
                        # Chuyển trang thành hình ảnh
                        pix = page.get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72))
                        img_path = os.path.join(temp_dir, f"page_{page_num+1}.png")
                        pix.save(img_path)

                        # OCR hình ảnh
                        try:
                            img = Image.open(img_path)
                            ocr_text = pytesseract.image_to_string(img, lang="vie+eng")
                            text = ocr_text
                        except Exception as ocr_err:
                            logging.warning(
                                f"OCR không thành công cho trang {page_num+1}: {str(ocr_err)}"
                            )

                    # Ghi văn bản vào tệp
                    if text.strip():
                        txt_file.write(f"--- Trang {page_num+1} ---\n\n")
                        txt_file.write(text)
                        txt_file.write("\n\n")

        # Kiểm tra kết quả
        with open(duong_dan_txt, "r", encoding="utf-8") as f:
            content = f.read().strip()

        if content:
            logging.info(
                f"Đã chuyển đổi thành công bằng PyMuPDF + OCR: '{duong_dan_txt}'"
            )
            return str(duong_dan_txt)
    except Exception as e:
        logging.warning(f"PyMuPDF + OCR không thành công: {str(e)}")

    # Phương pháp 3: Phương pháp dự phòng với PyPDF2
    try:
        logging.info("Đang thử phương pháp dự phòng với PyPDF2...")

        pdf = PdfReader(duong_dan_pdf)
        with open(duong_dan_txt, "w", encoding="utf-8") as txt_file:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    txt_file.write(f"--- Trang {i+1} ---\n\n")
                    txt_file.write(text)
                    txt_file.write("\n\n")

        logging.info(f"Đã chuyển đổi với phương pháp dự phòng: '{duong_dan_txt}'")
        return str(duong_dan_txt)
    except Exception as e:
        logging.error(f"Tất cả các phương pháp chuyển đổi đều thất bại: {str(e)}")

        # Tạo tệp trống nếu tất cả đều thất bại
        with open(duong_dan_txt, "w", encoding="utf-8") as f:
            f.write(f"Không thể trích xuất văn bản từ {Path(duong_dan_pdf).name}.\n")
            f.write(f"Lỗi: {str(e)}")

        return str(duong_dan_txt)


# Cập nhật BANG_CHUYEN_DOI với các hàm mới và cải tiến
BANG_CHUYEN_DOI: Dict[str, Dict[str, Callable]] = {
    ".pdf": {
        ".docx": chuyen_doi_pdf_sang_docx,
        ".xlsx": chuyen_doi_pdf_sang_xlsx,
        ".txt": chuyen_doi_pdf_sang_txt,
    },
    ".xlsx": {
        ".docx": chuyen_doi_xlsx_sang_docx,
        ".pdf": chuyen_doi_xlsx_sang_pdf,
        ".csv": chuyen_doi_xlsx_sang_csv,
    },
    ".docx": {
        ".pdf": chuyen_doi_docx_sang_pdf,
        ".xlsx": chuyen_doi_docx_sang_xlsx,
    },
    ".txt": {
        ".docx": chuyen_doi_txt_sang_docx,
        ".pdf": chuyen_doi_txt_sang_pdf,
        ".md": chuyen_doi_txt_sang_md,
    },
    ".csv": {
        ".xlsx": chuyen_doi_csv_sang_xlsx,
    },
    ".pptx": {
        ".pdf": chuyen_doi_pptx_sang_pdf,
        ".docx": chuyen_doi_pptx_sang_docx,
    },
    ".md": {
        ".txt": chuyen_doi_md_sang_txt,
        ".html": chuyen_doi_md_sang_html,
    },
}


def xu_ly_chuyen_doi(duong_dan_tep: str) -> None:
    duong_dan_tep = Path(duong_dan_tep)
    if not duong_dan_tep.is_file():
        logging.error(f"'{duong_dan_tep}' không phải là tệp hợp lệ hoặc không tồn tại.")
        return

    dinh_dang_nguon = duong_dan_tep.suffix.lower()
    if dinh_dang_nguon not in BANG_CHUYEN_DOI:
        logging.error(f"Định dạng tệp nguồn không được hỗ trợ: {dinh_dang_nguon}")
        return

    print(f"Các tùy chọn chuyển đổi có sẵn cho {dinh_dang_nguon}:")
    for i, dinh_dang_dich in enumerate(BANG_CHUYEN_DOI[dinh_dang_nguon].keys(), 1):
        print(f"{i}. {dinh_dang_nguon[1:].upper()} sang {dinh_dang_dich[1:].upper()}")

    lua_chon = input("Nhập số lựa chọn của bạn: ")
    try:
        lua_chon = int(lua_chon)
        dinh_dang_dich = list(BANG_CHUYEN_DOI[dinh_dang_nguon].keys())[lua_chon - 1]
    except (ValueError, IndexError):
        logging.error("Lựa chọn không hợp lệ.")
        return

    ham_chuyen_doi = BANG_CHUYEN_DOI[dinh_dang_nguon][dinh_dang_dich]
    ket_qua = ham_chuyen_doi(str(duong_dan_tep))

    if ket_qua:
        logging.info(f"Đã chuyển đổi thành công '{duong_dan_tep}' sang '{ket_qua}'")
    else:
        logging.error(f"Chuyển đổi thất bại cho '{duong_dan_tep}'")


def main():
    while True:
        print("\nChương trình chuyển đổi định dạng tệp")
        print("\nAuthor: tanbaycu")
        print("Vui lòng cung cấp đường dẫn tệp sử dụng \\ cho dấu phân cách thư mục.")
        duong_dan_tep = input("Nhập đường dẫn tệp (hoặc 'q' để thoát): ")

        if duong_dan_tep.lower() == "q":
            print("Cảm ơn bạn đã sử dụng chương trình. Tạm biệt!")
            break

        xu_ly_chuyen_doi(duong_dan_tep)

        lua_chon = input("\nBạn có muốn chuyển đổi tệp khác không? (y/n): ")
        if lua_chon.lower() != "y":
            print("Cảm ơn bạn đã sử dụng chương trình. Tạm biệt!")
            break


if __name__ == "__main__":
    main()
